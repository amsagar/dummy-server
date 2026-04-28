from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ────────────────────────────────────────────────────────────────────
NAVY    = RGBColor(0x08, 0x14, 0x33)   # deep navy – hero backgrounds
BLUE    = RGBColor(0x00, 0x52, 0xB4)   # primary blue – headers, cards
LBLUE   = RGBColor(0x00, 0x8C, 0xFF)   # accent bright blue – highlights
TEAL    = RGBColor(0x00, 0x7B, 0x83)   # teal – alternate accent
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY   = RGBColor(0xF3, 0xF6, 0xFA)   # light grey – slide backgrounds
MGRAY   = RGBColor(0xD0, 0xD8, 0xE8)   # mid grey – dividers
DGRAY   = RGBColor(0x3A, 0x48, 0x60)   # dark grey – body text
CHARCOAL= RGBColor(0x1A, 0x24, 0x3A)   # almost-black – contrast panels

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.50)
BLANK = prs.slide_layouts[6]


# ── helpers ────────────────────────────────────────────────────────────────────

def bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color


def rect(slide, l, t, w, h, fill, line_color=None, line_w=Pt(0)):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line_color:
        s.line.color.rgb = line_color
        s.line.width = line_w
    else:
        s.line.fill.background()
    return s


def tx(slide, text, l, t, w, h,
        size=Pt(13), bold=False, italic=False,
        color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = size
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color


def hbar(slide, t, h=0.07, color=LBLUE):
    rect(slide, 0, t, 13.33, h, color)


def slide_header(slide, title, subtitle=None, slide_num=None):
    rect(slide, 0, 0.07, 13.33, 1.05, NAVY)
    tx(slide, title, 0.42, 0.14, 9.50, 0.58,
       size=Pt(27), bold=True, color=WHITE)
    if subtitle:
        tx(slide, subtitle, 0.42, 0.72, 9.50, 0.35,
           size=Pt(11), italic=True, color=MGRAY)
    if slide_num:
        tx(slide, slide_num, 11.60, 0.22, 1.50, 0.38,
           size=Pt(10), color=MGRAY, align=PP_ALIGN.RIGHT)


def feature_row(slide, icon, title, body, t_pos):
    rect(slide, 0.35, t_pos, 7.30, 1.05, WHITE, BLUE, Pt(0.7))
    rect(slide, 0.35, t_pos, 0.58, 1.05, BLUE)
    tx(slide, icon, 0.35, t_pos + 0.20, 0.58, 0.60,
       size=Pt(20), align=PP_ALIGN.CENTER, color=WHITE)
    tx(slide, title, 1.02, t_pos + 0.07, 6.50, 0.36,
       size=Pt(12), bold=True, color=NAVY)
    tx(slide, body, 1.02, t_pos + 0.44, 6.50, 0.54,
       size=Pt(10.5), color=DGRAY)


def uc_card(slide, l, t, icon, title, lines):
    w, h = 2.95, 5.10
    rect(slide, l, t, w, h, WHITE, BLUE, Pt(1.0))
    rect(slide, l, t, w, 0.90, BLUE)
    tx(slide, icon, l, t + 0.04, w, 0.50,
       size=Pt(22), align=PP_ALIGN.CENTER, color=WHITE)
    tx(slide, title, l + 0.10, t + 0.56, w - 0.20, 0.34,
       size=Pt(10.5), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb = slide.shapes.add_textbox(
        Inches(l + 0.15), Inches(t + 1.00), Inches(w - 0.30), Inches(h - 1.10))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(5)
        r = p.add_run()
        r.text = line
        r.font.size = Pt(10)
        r.font.color.rgb = DGRAY


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Hero
# ══════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)
bg(s1, NAVY)

# left colour panel
rect(s1, 0, 0, 5.60, 7.50, BLUE)

# vertical bright accent stripe
rect(s1, 5.60, 0, 0.14, 7.50, LBLUE)

# top-left label
tx(s1, "ENTERPRISE AI", 0.38, 0.28, 4.80, 0.48,
   size=Pt(13), bold=True, color=LBLUE, italic=True)

# main headline
tx(s1, "Intelligent\nAgent Platform", 0.38, 0.82, 4.90, 2.10,
   size=Pt(38), bold=True, color=WHITE)

# sub-headline
tx(s1, "Multi-model orchestration · Streaming chat · Human-in-the-loop\n"
       "Smart tool routing · Persistent memory · Audit-grade logging",
   0.38, 3.05, 4.90, 1.00,
   size=Pt(12), color=RGBColor(0xBB, 0xCC, 0xEE))

# right-side value statement
tx(s1, "Connect any API.\nAutomate any workflow.\nKeep humans in control.",
   6.10, 2.20, 6.80, 2.20,
   size=Pt(26), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

tx(s1, "Powered by Spring AI  ·  pgvector  ·  React",
   6.20, 4.55, 6.40, 0.45,
   size=Pt(12), color=LBLUE, align=PP_ALIGN.CENTER, italic=True)

# capability pills (bottom-left)
pills = ["OpenAI", "Azure OpenAI", "Anthropic Claude", "Ollama", "Google Vertex"]
pill_x = 0.38
for p_label in pills:
    pw = 1.10
    rect(s1, pill_x, 4.70, pw, 0.35, CHARCOAL)
    tx(s1, p_label, pill_x, 4.72, pw, 0.30,
       size=Pt(9), color=LGRAY, align=PP_ALIGN.CENTER, bold=True)
    pill_x += pw + 0.08

# bottom footer bar
hbar(s1, 7.22, 0.28, LBLUE)
tx(s1, "Confidential  |  Client Presentation  |  2026",
   0.30, 7.22, 12.70, 0.26,
   size=Pt(9), color=NAVY, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Use Cases
# ══════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)
bg(s2, LGRAY)
hbar(s2, 0.0, 0.07, LBLUE)
hbar(s2, 7.43, 0.07, BLUE)

slide_header(s2,
             "Real-World Use Cases",
             "What the platform does — out of the box",
             "2 / 4")

use_cases = [
    ("⚙️", "Automated\nWorkflow Orchestration",
     ["• Chain multi-step operations across APIs",
      "• LLM decides which tools to call & in what order",
      "• Built-in retry, timeout, and error handling",
      "• Approval gates before high-risk actions",
      "• Full audit log of every tool call"]),
    ("🔌", "System Integration\n& API Automation",
     ["• Import OpenAPI specs in one click",
      "• Convert cURL commands to live tools",
      "• MCP server discovery & tool routing",
      "• OAuth2, API Key, Bearer, Basic auth",
      "• Works with any REST or JSON endpoint"]),
    ("🧠", "Knowledge-Driven\nIntelligent Validation",
     ["• Encode domain rules as Markdown skills",
      "• Semantic embedding retrieves right rules",
      "• LLM cross-references data against rules",
      "• Structured PASS / FLAG / CRITICAL reports",
      "• Rules versioned and auditable"]),
    ("🛡️", "Human-in-the-Loop\nCompliance & Oversight",
     ["• Configurable approval gates per tool",
      "• Multi-choice interaction prompts",
      "• Real-time streaming: see reasoning live",
      "• Every decision persisted for compliance",
      "• Role-based access per tool scope"]),
]

x_positions = [0.35, 3.42, 6.49, 9.56]
for (icon, title, lines), x in zip(use_cases, x_positions):
    uc_card(s2, x, 1.27, icon, title, lines)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Features & Architecture
# ══════════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK)
bg(s3, LGRAY)
hbar(s3, 0.0, 0.07, BLUE)
hbar(s3, 7.43, 0.07, LBLUE)

slide_header(s3,
             "Key Features & Architecture",
             "What makes this platform enterprise-ready",
             "3 / 4")

features = [
    ("⚡", "Real-Time Streaming Chat",
     "SSE-based streaming with live reasoning visibility. Supports extended thinking, "
     "cancellation, and partial-response recovery."),
    ("🔍", "Semantic Tool Selection",
     "pgvector embeddings rank the most relevant tools per query. Usage signals "
     "boost frequently successful tools. Cap of 48 tools per turn avoids prompt bloat."),
    ("📝", "Persistent Memory System",
     "User-scoped memories stored in Azure Blob. LLM reads/writes structured Markdown "
     "memory files across sessions. Fully auditable."),
    ("🧩", "Skills & Domain Knowledge",
     "Markdown-based skill files encode business rules, prompts, and validation logic. "
     "Lazy-loaded — catalog in system prompt, content fetched on demand."),
    ("💰", "Cost & Usage Tracking",
     "Per-turn, per-session, per-model token accounting. Budget alerts at configurable "
     "thresholds. Prompt caching hints for Anthropic models."),
]

for i, (icon, title, body) in enumerate(features):
    feature_row(s3, icon, title, body, 1.27 + i * 1.22)

# right panel — architecture stack
rect(s3, 7.95, 1.22, 5.05, 6.00, WHITE, NAVY, Pt(1.0))
rect(s3, 7.95, 1.22, 5.05, 0.44, NAVY)
tx(s3, "Platform Architecture", 7.95, 1.26, 5.05, 0.36,
   size=Pt(12), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

arch = [
    (BLUE,     "React UI  +  Streaming Chat"),
    (CHARCOAL, "REST / SSE  (Spring MVC)"),
    (NAVY,     "Agent Orchestrator  (Java 21)"),
    (CHARCOAL, "Spring AI  |  MCP Client"),
    (BLUE,     "LLM Router  —  Multi-Model"),
    (CHARCOAL, "Tool Callbacks  |  HITL Gates"),
    (TEAL,     "Tool Registry  +  Skills Engine"),
    (CHARCOAL, "pgvector  |  Embedding Index"),
    (TEAL,     "PostgreSQL  +  Azure Blob"),
    (CHARCOAL, "OpenAPI  ·  cURL  ·  MCP Tools"),
]

for i, (col, label) in enumerate(arch):
    tp = 1.76 + i * 0.54
    rect(s3, 8.15, tp, 4.65, 0.46, col)
    tx(s3, label, 8.15, tp + 0.06, 4.65, 0.36,
       size=Pt(10), bold=(col != CHARCOAL), color=WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Business Value & Differentiators
# ══════════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK)
bg(s4, NAVY)
hbar(s4, 0.0, 0.07, LBLUE)

rect(s4, 0, 0.07, 13.33, 1.05, CHARCOAL)
tx(s4, "Business Value & Differentiators", 0.42, 0.13, 9.50, 0.58,
   size=Pt(27), bold=True, color=WHITE)
tx(s4, "4 / 4", 11.60, 0.22, 1.50, 0.38,
   size=Pt(10), color=MGRAY, align=PP_ALIGN.RIGHT)

# metric tiles
metrics = [
    ("Any LLM",    "OpenAI · Azure · Anthropic\nOllama · Google Vertex"),
    ("No-Code",    "OpenAPI & cURL import\ntools in one click"),
    ("< 2 min",    "Full multi-step workflow\nfrom prompt to result"),
    ("100%",       "Tool calls logged\nand auditable"),
]
for i, (num, label) in enumerate(metrics):
    l = 0.38 + i * 3.23
    rect(s4, l, 1.28, 2.95, 1.55, BLUE)
    rect(s4, l, 1.28, 2.95, 0.07, LBLUE)
    tx(s4, num, l, 1.40, 2.95, 0.78,
       size=Pt(30), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tx(s4, label, l, 2.22, 2.95, 0.56,
       size=Pt(10), color=RGBColor(0xBB, 0xCC, 0xFF), align=PP_ALIGN.CENTER)

# differentiator bullets (left column)
tx(s4, "Why this platform stands out", 0.42, 3.05, 6.40, 0.40,
   size=Pt(13), bold=True, color=LBLUE)

diffs = [
    "Provider-agnostic — swap LLMs without re-engineering any workflow",
    "Skill-first design — domain experts write rules in Markdown, not code",
    "HITL built-in — approval gates, not afterthoughts, on every risky action",
    "Semantic tool routing — the right API is called every time, automatically",
    "Enterprise security — JWT auth, encrypted credentials, scope enforcement",
    "Full observability — every token, tool call, and decision is persisted",
]
for i, d in enumerate(diffs):
    tx(s4, f"  {d}", 0.42, 3.52 + i * 0.54, 6.60, 0.50,
       size=Pt(11), color=WHITE)
    rect(s4, 0.42, 3.56 + i * 0.54, 0.08, 0.32, LBLUE)

# right column — deployment highlights
rect(s4, 7.20, 3.00, 5.80, 3.90, CHARCOAL)
rect(s4, 7.20, 3.00, 5.80, 0.44, BLUE)
tx(s4, "Deployment & Integration", 7.20, 3.04, 5.80, 0.36,
   size=Pt(12), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

deploy = [
    ("Backend",     "Spring Boot 3 · Java 21 · Maven"),
    ("Database",    "PostgreSQL + pgvector extension"),
    ("Storage",     "Azure Blob Storage (skills & memory)"),
    ("Frontend",    "React 18 · TypeScript · Tailwind"),
    ("Auth",        "JWT · OAuth2 · API Key · Basic"),
    ("LLM Access",  "Spring AI  — all major providers"),
    ("Protocols",   "REST · SSE · MCP · OpenAPI 3.0"),
]
for i, (label, val) in enumerate(deploy):
    tp = 3.52 + i * 0.50
    tx(s4, label, 7.32, tp, 1.60, 0.44,
       size=Pt(10), bold=True, color=LBLUE)
    tx(s4, val, 8.96, tp, 3.90, 0.44,
       size=Pt(10), color=WHITE)

# bottom footer
hbar(s4, 7.22, 0.28, LBLUE)
tx(s4, "Confidential  |  Client Presentation  |  2026",
   0.30, 7.22, 12.70, 0.26,
   size=Pt(9), color=NAVY, align=PP_ALIGN.RIGHT)


# ── save ───────────────────────────────────────────────────────────────────────
import os
out = os.path.join(os.path.expanduser("~"), "Desktop", "AI_Agent_Platform_Presentation.pptx")
prs.save(out)
print("Saved:", out)
