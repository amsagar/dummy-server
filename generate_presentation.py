from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree

# ── Brand colours ──────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x0D, 0x1B, 0x3E)   # deep navy – backgrounds
BLUE   = RGBColor(0x1A, 0x6B, 0xC4)   # PODS blue – accents
ORANGE = RGBColor(0xF4, 0x7B, 0x20)   # PODS orange – highlights
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xF0, 0xF4, 0xF8)   # light grey – card backgrounds
DGRAY  = RGBColor(0x44, 0x55, 0x66)   # dark grey – body text

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # completely blank layout

# ── helpers ────────────────────────────────────────────────────────────────────

def bg(slide, color):
    """Fill slide background with a solid colour."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, l, t, w, h, fill_color, line_color=None, line_width=Pt(0)):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def txbox(slide, text, l, t, w, h,
          font_size=Pt(14), bold=False, color=WHITE,
          align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def accent_bar(slide, t=0.0, h=0.07, color=ORANGE):
    rect(slide, 0, t, 13.33, h, color)

def card(slide, l, t, w, h, title, body_lines, icon=""):
    """Rounded-corner card with title + bullet body."""
    r = rect(slide, l, t, w, h, LGRAY, BLUE, Pt(1.2))
    # title strip
    rect(slide, l, t, w, 0.42, BLUE)
    txbox(slide, f"{icon}  {title}" if icon else title,
          l+0.12, t+0.04, w-0.24, 0.36,
          font_size=Pt(12), bold=True, color=WHITE)
    # body
    body_tb = slide.shapes.add_textbox(
        Inches(l+0.15), Inches(t+0.50), Inches(w-0.30), Inches(h-0.60))
    tf = body_tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(body_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(10.5)
        run.font.color.rgb = DGRAY

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title / Hero
# ══════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)
bg(s1, NAVY)

# left colour panel
rect(s1, 0, 0, 5.5, 7.5, BLUE)

# orange accent stripe
rect(s1, 5.5, 0, 0.12, 7.5, ORANGE)

# PODS logo text (top-left)
txbox(s1, "PODS", 0.35, 0.30, 2.5, 0.65,
      font_size=Pt(32), bold=True, color=WHITE)
txbox(s1, "Portable On Demand Storage", 0.35, 0.90, 4.6, 0.40,
      font_size=Pt(11), color=RGBColor(0xAA, 0xCC, 0xFF), italic=True)

# main headline
txbox(s1, "AI-Powered\nOrder Validation\nAgent", 0.35, 2.00, 4.8, 2.20,
      font_size=Pt(34), bold=True, color=WHITE)

# sub-headline
txbox(s1, "Automated multi-system validation across\nSalesforce · D365 · RBMS · POET · PodHunter",
      0.35, 4.30, 4.8, 0.90,
      font_size=Pt(13), color=RGBColor(0xCC, 0xDD, 0xFF))

# right-side tagline block
txbox(s1, "Catch billing errors\nbefore they cost you.",
      6.00, 2.60, 6.80, 1.20,
      font_size=Pt(26), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txbox(s1, "Real-time · Intelligent · Actionable",
      6.20, 3.90, 6.40, 0.50,
      font_size=Pt(14), color=ORANGE, align=PP_ALIGN.CENTER, italic=True)

# bottom bar
accent_bar(s1, t=7.20, h=0.30, color=ORANGE)
txbox(s1, "Confidential — Client Presentation  |  2026",
      0.30, 7.20, 12.70, 0.28,
      font_size=Pt(9), color=WHITE, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — The Problem & Use Cases
# ══════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)
bg(s2, WHITE)
accent_bar(s2, t=0.0, h=0.08, color=ORANGE)
accent_bar(s2, t=7.42, h=0.08, color=BLUE)

# header band
rect(s2, 0, 0.08, 13.33, 1.10, NAVY)
txbox(s2, "The Problem  →  The Use Cases", 0.40, 0.14, 9.0, 0.55,
      font_size=Pt(26), bold=True, color=WHITE)
txbox(s2, "Slide 2 of 4", 11.50, 0.22, 1.60, 0.40,
      font_size=Pt(10), color=RGBColor(0xAA, 0xBB, 0xCC), align=PP_ALIGN.RIGHT)

# problem statement
txbox(s2,
      "PODS manages thousands of active container rentals across multiple disconnected systems. "
      "A single order touches Salesforce CPQ, D365 F&O, RBMS, POET, and PodHunter — "
      "and any gap between them can silently generate overbilling, missed revenue, or blocked invoicing.",
      0.40, 1.30, 12.50, 0.80,
      font_size=Pt(12), color=DGRAY)

# 4 use-case cards
uc = [
    ("🔍", "Billing Integrity Check",
     ["Detect post-FPU auto-renew leakage",
      "Catch missed rent generation gaps",
      "Flag double-billing (overlapping lines)",
      "Monitor 48-hr billing cliff window"]),
    ("🔗", "D365 Sync Validation",
     ["Verify 35 D365 validation rules",
      "Detect Booked → Invoiced stuck lines",
      "Check FPU propagation to ERP",
      "Identify blank Line Identity failures"]),
    ("📦", "Container & Ghost Detection",
     ["Cross-check PodHunter vs Salesforce CID",
      "Detect ghost assignments post-FPU",
      "Find orphaned legs & order lines",
      "Validate container on-rent status"]),
    ("🗓️", "Service Leg & Timestamp Audit",
     ["Validate IDEL→WRT→RDL→FPU patterns",
      "Detect out-of-sequence timestamps",
      "Check POET propagation to D365",
      "Flag retroactive timestamp corrections"]),
]

positions = [(0.35, 2.25), (3.55, 2.25), (6.75, 2.25), (9.95, 2.25)]
for (icon, title, body), (l, t) in zip(uc, positions):
    card(s2, l, t, 3.00, 4.90, title, body, icon)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Key Features & Architecture
# ══════════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK)
bg(s3, LGRAY)
accent_bar(s3, t=0.0, h=0.08, color=BLUE)
accent_bar(s3, t=7.42, h=0.08, color=ORANGE)

rect(s3, 0, 0.08, 13.33, 1.10, NAVY)
txbox(s3, "Key Features & Architecture", 0.40, 0.14, 9.0, 0.55,
      font_size=Pt(26), bold=True, color=WHITE)
txbox(s3, "Slide 3 of 4", 11.50, 0.22, 1.60, 0.40,
      font_size=Pt(10), color=RGBColor(0xAA, 0xBB, 0xCC), align=PP_ALIGN.RIGHT)

# ── left column: feature list ──────────────────────────────────────────────
features = [
    ("⚡", "Real-Time Streaming Agent",
     "Spring AI orchestrator streams responses via SSE. Supports Claude, GPT-4, Gemini and any OpenAI-compatible model."),
    ("🧠", "Skill-Based Knowledge Engine",
     "Domain rules (billing, D365, ghost detection, timestamps) loaded as vector-embedded skills — retrieved at query time."),
    ("🔒", "Multi-Tenant & Secure",
     "JWT auth, per-tenant model routing, role-based tool access, and configurable CORS policies."),
    ("🔌", "MCP Tool Integration",
     "Model Context Protocol support lets the agent call live APIs across all 5 PODS systems in a single conversation turn."),
    ("📋", "Structured Validation Reports",
     "Every investigation produces a standardised 6-check report: PASS / FLAG / CRITICAL with team assignments and deadlines."),
]

for i, (icon, title, desc) in enumerate(features):
    t_pos = 1.35 + i * 1.14
    rect(s3, 0.35, t_pos, 7.20, 1.02, WHITE, BLUE, Pt(0.8))
    rect(s3, 0.35, t_pos, 0.55, 1.02, BLUE)
    txbox(s3, icon, 0.35, t_pos+0.22, 0.55, 0.55,
          font_size=Pt(18), align=PP_ALIGN.CENTER, color=WHITE)
    txbox(s3, title, 1.00, t_pos+0.06, 6.40, 0.36,
          font_size=Pt(12), bold=True, color=NAVY)
    txbox(s3, desc, 1.00, t_pos+0.44, 6.40, 0.52,
          font_size=Pt(10), color=DGRAY)

# ── right column: architecture diagram (text-based) ───────────────────────
rect(s3, 7.90, 1.30, 5.10, 5.90, WHITE, NAVY, Pt(1.2))
txbox(s3, "System Architecture", 7.90, 1.30, 5.10, 0.42,
      font_size=Pt(12), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
rect(s3, 7.90, 1.30, 5.10, 0.42, NAVY)
txbox(s3, "System Architecture", 7.90, 1.34, 5.10, 0.36,
      font_size=Pt(12), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

arch_items = [
    (BLUE,   "React UI  +  Chat Interface"),
    (DGRAY,  "↕  REST / SSE"),
    (NAVY,   "Spring Boot Agent  (Java 21)"),
    (DGRAY,  "↕  Spring AI  |  MCP Client"),
    (BLUE,   "LLM Router  (Claude / GPT / Gemini)"),
    (DGRAY,  "↕  Tool Callbacks"),
    (ORANGE, "Salesforce CPQ  ·  D365 F&O"),
    (ORANGE, "RBMS  ·  POET  ·  PodHunter"),
    (DGRAY,  "↕  pgvector"),
    (BLUE,   "Skills Knowledge Base"),
]

for i, (col, label) in enumerate(arch_items):
    t_pos = 1.82 + i * 0.50
    rect(s3, 8.10, t_pos, 4.70, 0.40, col)
    txbox(s3, label, 8.10, t_pos+0.04, 4.70, 0.34,
          font_size=Pt(10), bold=(col != DGRAY), color=WHITE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Value & Next Steps
# ══════════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK)
bg(s4, NAVY)
accent_bar(s4, t=0.0, h=0.08, color=ORANGE)

rect(s4, 0, 0.08, 13.33, 1.10, RGBColor(0x0A, 0x14, 0x30))
txbox(s4, "Business Value & Next Steps", 0.40, 0.14, 9.0, 0.55,
      font_size=Pt(26), bold=True, color=WHITE)
txbox(s4, "Slide 4 of 4", 11.50, 0.22, 1.60, 0.40,
      font_size=Pt(10), color=RGBColor(0x88, 0x99, 0xBB), align=PP_ALIGN.RIGHT)

# ── value metrics ──────────────────────────────────────────────────────────
metrics = [
    ("$0", "Overbilling leakage\nwith auto-renew guard"),
    ("6", "Validation checks\nper order, automated"),
    ("35", "D365 rules enforced\nbefore sync failure"),
    ("<2 min", "Full order audit\nend-to-end"),
]

for i, (num, label) in enumerate(metrics):
    l = 0.40 + i * 3.22
    rect(s4, l, 1.35, 2.90, 1.60, BLUE)
    rect(s4, l, 1.35, 2.90, 0.08, ORANGE)
    txbox(s4, num, l, 1.50, 2.90, 0.80,
          font_size=Pt(34), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txbox(s4, label, l, 2.28, 2.90, 0.60,
          font_size=Pt(10), color=RGBColor(0xCC, 0xDD, 0xFF), align=PP_ALIGN.CENTER)

# ── benefits ───────────────────────────────────────────────────────────────
benefits = [
    "✅  Eliminate silent overbilling — auto-renew leakage caught before the next anniversary date",
    "✅  Prevent revenue gaps — continuous rent generation verified from IDEL through FPU",
    "✅  Reduce manual investigation time — structured 6-check report replaces ad-hoc queries",
    "✅  Proactive billing cliff alerts — 48-hour window monitoring flags at-risk containers",
    "✅  Cross-system data integrity — ghost assignments and orphaned records surfaced automatically",
]

for i, b in enumerate(benefits):
    txbox(s4, b, 0.50, 3.20 + i * 0.52, 12.30, 0.46,
          font_size=Pt(11.5), color=WHITE)

# ── next steps ─────────────────────────────────────────────────────────────
rect(s4, 0.35, 6.00, 12.60, 1.10, RGBColor(0x0A, 0x14, 0x30))
rect(s4, 0.35, 6.00, 12.60, 0.06, ORANGE)

steps = [
    "Phase 1\nRead-Only Validation",
    "Phase 2\nWrite-Back Corrections",
    "Phase 3\nProactive Monitoring",
    "Phase 4\nFull Ops Integration",
]
for i, step in enumerate(steps):
    l = 0.55 + i * 3.15
    rect(s4, l, 6.10, 2.80, 0.85, BLUE)
    txbox(s4, step, l, 6.14, 2.80, 0.78,
          font_size=Pt(10), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

accent_bar(s4, t=7.20, h=0.30, color=ORANGE)
txbox(s4, "Confidential — Client Presentation  |  PODS Order Validation Agent  |  2026",
      0.30, 7.20, 12.70, 0.28,
      font_size=Pt(9), color=WHITE, align=PP_ALIGN.RIGHT)

# ── save ───────────────────────────────────────────────────────────────────
import os
out = os.path.join(os.path.expanduser("~"), "Desktop", "PODS_Agent_Presentation.pptx")
prs.save(out)
with open("/tmp/ppt_done.txt", "w") as f:
    f.write("saved:" + out + "\n")
